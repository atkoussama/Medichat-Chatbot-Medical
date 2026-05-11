import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import fitz  # PyMuPDF for PDF extraction
from pptx.dml.color import RGBColor

# --- Helper functions ---
def extract_pptx_titles(pptx_path, max_slides=5):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        if i >= max_slides:
            break
        title = ''
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not title:
                    title = text
                else:
                    content.append(text)
        slides.append({'title': title, 'content': content})
    return slides

def extract_pdf_headings(pdf_path, max_pages=3):
    doc = fitz.open(pdf_path)
    structure = []
    for page_num in range(min(max_pages, doc.page_count)):
        page = doc.load_page(page_num)
        text = page.get_text()
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        if lines:
            structure.append({'page': page_num+1, 'lines': lines[:8]})
    return structure

def get_medi_chat_content():
    # Summarized content based on README and .tex analysis
    return [
        {
            'title': 'Introduction',
            'content': [
                'MediChat is an AI-powered medical assistant web app.',
                'It provides users with preliminary medical information and guidance.',
                'Developed using React, Tailwind CSS, and AI integration.'
            ]
        },
        {
            'title': 'Objectives',
            'content': [
                '- Intuitive UI for medical consultations',
                '- Multilingual and accessible',
                '- Voice interaction (TTS/STT)',
                '- Responsive and modern design',
                '- User privacy and data security'
            ]
        },
        {
            'title': 'Key Features',
            'content': [
                '• Multilingual (English, French, Arabic, RTL support)',
                '• Voice interaction (text-to-speech, speech-to-text)',
                '• 3D doctor avatar and animated UI',
                '• Context-aware AI responses',
                '• Profile management, news, and accessibility options'
            ]
        },
        {
            'title': 'System Architecture',
            'content': [
                '1. User Interface: React, Tailwind CSS',
                '2. Service Layer: Chat, voice, and API integration',
                '3. External APIs: Mistral AI, Web Speech API'
            ]
        },
        {
            'title': 'UI/UX Design',
            'content': [
                '• Modern, clean, and responsive design',
                '• Glassmorphism effects',
                '• High contrast and accessibility',
                '• Animated message bubbles and avatars'
            ]
        },
        {
            'title': 'Technologies Used',
            'content': [
                '• React.js, Tailwind CSS, React Router',
                '• React Three Fiber (3D avatar)',
                '• Web Speech API',
                '• Axios, OpenRouter API',
                '• Component-based architecture'
            ]
        },
        {
            'title': 'Challenges & Solutions',
            'content': [
                '• Speech API browser compatibility: custom fallback',
                '• 3D avatar performance: lazy loading, 2D fallback',
                '• Conversation context: message history system'
            ]
        },
        {
            'title': 'Conclusion',
            'content': [
                'MediChat bridges the gap between users and medical information.',
                'Combines AI, modern web tech, and accessibility for digital health.',
                'Ready for future enhancements and broader impact.'
            ]
        }
    ]

def add_theme_background(slide, prs):
    # Use presentation's slide dimensions
    left = 0
    top = prs.slide_height - Inches(1.2)
    width = prs.slide_width
    height = Inches(1.2)
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(220, 235, 250)  # Light blue
    line = shape.line
    line.color.rgb = RGBColor(180, 200, 230)
    # Add a gear/circle in the corner (simple circle as placeholder)
    circle = slide.shapes.add_shape(
        9,  # Oval
        prs.slide_width - Inches(1.2), Inches(0.1), Inches(1), Inches(1)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(180, 200, 230)
    circle.line.color.rgb = RGBColor(120, 150, 200)

def add_architecture_diagram(slide):
    left = Inches(1)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1)
    # UI box
    ui = slide.shapes.add_shape(1, left, top, width, height)
    ui.text = "User Interface\n(React, Tailwind)"
    ui.fill.solid()
    ui.fill.fore_color.rgb = RGBColor(100, 180, 255)
    # Service box
    service_left = left + Inches(2.5)
    service = slide.shapes.add_shape(1, service_left, top, width, height)
    service.text = "Service Layer\n(Chat, Voice)"
    service.fill.solid()
    service.fill.fore_color.rgb = RGBColor(120, 200, 255)
    # API box
    api_left = left + Inches(5)
    api = slide.shapes.add_shape(1, api_left, top, width, height)
    api.text = "External APIs\n(Mistral AI, Web Speech)"
    api.fill.solid()
    api.fill.fore_color.rgb = RGBColor(140, 220, 255)
    # Arrows (connectors)
    # UI -> Service
    start_x = left + width
    start_y = top + height / 2
    end_x = service_left
    end_y = top + height / 2
    slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    # Service -> API
    start_x2 = service_left + width
    start_y2 = top + height / 2
    end_x2 = api_left
    end_y2 = top + height / 2
    slide.shapes.add_connector(1, start_x2, start_y2, end_x2, end_y2)

def add_user_journey_flow(slide):
    left = Inches(1)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1)
    login = slide.shapes.add_shape(1, left, top, width, height)
    login.text = "Login/Register"
    login.fill.solid()
    login.fill.fore_color.rgb = RGBColor(180, 220, 255)
    chat_left = left + Inches(2.5)
    chat = slide.shapes.add_shape(1, chat_left, top, width, height)
    chat.text = "Chat with AI"
    chat.fill.solid()
    chat.fill.fore_color.rgb = RGBColor(120, 180, 255)
    result_left = left + Inches(5)
    result = slide.shapes.add_shape(1, result_left, top, width, height)
    result.text = "Get Medical Info"
    result.fill.solid()
    result.fill.fore_color.rgb = RGBColor(80, 140, 220)
    # Arrows (connectors)
    # Login -> Chat
    start_x = left + width
    start_y = top + height / 2
    end_x = chat_left
    end_y = top + height / 2
    slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    # Chat -> Result
    start_x2 = chat_left + width
    start_y2 = top + height / 2
    end_x2 = result_left
    end_y2 = top + height / 2
    slide.shapes.add_connector(1, start_x2, start_y2, end_x2, end_y2)

def generate_medi_chat_pptx(output_path, medi_chat_slides, template_slides=None):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "MediChat: Medical Chatbot Project"
    slide.placeholders[1].text = "AI-Powered Medical Assistant\nFinal Project Presentation"
    add_theme_background(slide, prs)
    # Add MediChat slides
    for s in medi_chat_slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s['title']
        slide.placeholders[1].text = '\n'.join(s['content'])
        add_theme_background(slide, prs)
        # Add diagrams to relevant slides
        if s['title'] == 'System Architecture':
            add_architecture_diagram(slide)
        if s['title'] == 'User Journey' or s['title'] == 'Key Features':
            add_user_journey_flow(slide)
    # Optionally add a few template-inspired slides at the end
    if template_slides:
        for tslide in template_slides:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = tslide['title'][:60]
            content = '\n'.join(tslide['content'])
            slide.placeholders[1].text = content[:500]
            add_theme_background(slide, prs)
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    pptx_path = "PRESENTATION PFA FINALE..pptx"
    pdf_paths = ["CHATBOT ENSIAS.pdf", "breast cancer prediction.pdf"]
    output_path = "MediChat_Presentation.pptx"

    # Extract a few template slides for inspiration
    template_slides = extract_pptx_titles(pptx_path, max_slides=2)
    # Optionally, extract from PDFs (not used directly for slides, but for inspiration)
    pdf_structures = [extract_pdf_headings(pdf) for pdf in pdf_paths]
    # Get MediChat content
    medi_chat_slides = get_medi_chat_content()
    # Generate new PPTX
    generate_medi_chat_pptx(output_path, medi_chat_slides, template_slides)

    print("\nDone! Your MediChat presentation is ready as MediChat_Presentation.pptx.")
    print("\nIf you need to install dependencies, run:")
    print("pip install python-pptx pymupdf")
